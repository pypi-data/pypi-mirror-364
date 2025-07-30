"""
PPTX parser implementation using python-pptx
"""

import asyncio
import logging
from typing import List, Dict, Any, TYPE_CHECKING
from pathlib import Path
import io

try:
    from pptx import Presentation

    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    # Create dummy class for type hints
    if TYPE_CHECKING:
        from pptx import Presentation
    else:
        Presentation = object

from .base import BaseParser
from ..core.models import ParsedDocument, ParserConfig, FileType, ContentBlock
from ..core.exceptions import ProcessingError


class PPTXParser(BaseParser):
    """Parser for PPTX documents"""

    def __init__(self):
        super().__init__()
        self.supported_formats = [FileType.PPTX]

        if not HAS_PPTX:
            raise ImportError(
                "PPTX parsing requires python-pptx. Install with: pip install python-pptx"
            )

    async def parse_async(
        self, file_path: Path, config: ParserConfig
    ) -> ParsedDocument:
        """Parse PPTX file asynchronously"""
        try:
            loop = asyncio.get_event_loop()
            return await loop.run_in_executor(
                None, self._parse_pptx_sync, str(file_path), config
            )
        except Exception as e:
            raise ProcessingError(f"Failed to parse PPTX: {str(e)}", e)

    async def parse_from_bytes_async(
        self, data: bytes, filename: str, config: ParserConfig
    ) -> ParsedDocument:
        """Parse PPTX from bytes asynchronously"""
        try:
            loop = asyncio.get_event_loop()
            return await loop.run_in_executor(
                None, self._parse_pptx_from_bytes_sync, data, filename, config
            )
        except Exception as e:
            raise ProcessingError(f"Failed to parse PPTX from bytes: {str(e)}", e)

    def _parse_pptx_sync(self, file_path: str, config: ParserConfig) -> ParsedDocument:
        """Synchronous PPTX parsing"""
        path = Path(file_path)
        document = asyncio.run(
            self._create_base_document(
                path, path.name, FileType.PPTX, path.stat().st_size
            )
        )

        try:
            prs = Presentation(file_path)
            return self._extract_content(prs, document, config)
        except Exception as e:
            raise ProcessingError(f"Failed to parse PPTX: {str(e)}", e)

    def _parse_pptx_from_bytes_sync(
        self, data: bytes, filename: str, config: ParserConfig
    ) -> ParsedDocument:
        """Synchronous PPTX parsing from bytes"""
        document = asyncio.run(
            self._create_base_document(None, filename, FileType.PPTX, len(data))
        )

        try:
            prs = Presentation(io.BytesIO(data))
            return self._extract_content(prs, document, config)
        except Exception as e:
            raise ProcessingError(f"Failed to parse PPTX from bytes: {str(e)}", e)

    def _extract_content(
        self, prs, document: ParsedDocument, config: ParserConfig
    ) -> ParsedDocument:
        """Extract content from PPTX presentation"""
        """Extract content from PPTX presentation"""
        content_parts = []
        content_blocks = []
        tables = []
        images = []
        links = []

        # Extract core properties
        if prs.core_properties:
            props = prs.core_properties
            document.metadata.title = props.title
            document.metadata.author = props.author
            document.metadata.subject = props.subject
            document.metadata.creation_date = props.created
            document.metadata.modification_date = props.modified

        # Update page count (slides)
        document.metadata.page_count = len(prs.slides)

        for slide_idx, slide in enumerate(prs.slides):
            slide_content = []

            # Extract slide title if present
            if hasattr(slide, "shapes"):
                title_shape = None
                for shape in slide.shapes:
                    if (
                        hasattr(shape, "placeholder_format")
                        and shape.placeholder_format
                    ):
                        if shape.placeholder_format.type == 1:  # Title placeholder
                            title_shape = shape
                            break

                if title_shape and hasattr(title_shape, "text"):
                    title_text = title_shape.text.strip()
                    if title_text:
                        content_blocks.append(
                            ContentBlock(
                                content=title_text,
                                block_type="title",
                                page_number=slide_idx + 1,
                            )
                        )
                        slide_content.append(f"Title: {title_text}")

            # Extract text from all shapes
            for shape in slide.shapes:
                try:
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text and text not in slide_content:  # Avoid duplicates
                            block_type = "text"

                            # Detect if this is a title or subtitle
                            if (
                                hasattr(shape, "placeholder_format")
                                and shape.placeholder_format
                            ):
                                if shape.placeholder_format.type == 1:
                                    block_type = "title"
                                elif shape.placeholder_format.type == 2:
                                    block_type = "subtitle"

                            content_blocks.append(
                                ContentBlock(
                                    content=text,
                                    block_type=block_type,
                                    page_number=slide_idx + 1,
                                )
                            )
                            slide_content.append(text)

                    # Extract tables if enabled
                    if config.extract_tables and hasattr(shape, "table"):
                        try:
                            table_data = []
                            for row in shape.table.rows:
                                row_data = []
                                for cell in row.cells:
                                    row_data.append(cell.text.strip())
                                table_data.append(row_data)

                            if table_data:
                                tables.append(
                                    {"slide": slide_idx + 1, "data": table_data}
                                )

                                # Add table as content block
                                table_text = "\n".join(
                                    ["\t".join(row) for row in table_data]
                                )
                                content_blocks.append(
                                    ContentBlock(
                                        content=table_text,
                                        block_type="table",
                                        page_number=slide_idx + 1,
                                    )
                                )
                                slide_content.append(table_text)
                        except Exception as e:
                            logging.warning(
                                f"Failed to extract table from slide {slide_idx + 1}: {str(e)}"
                            )

                    # Extract images if enabled
                    if config.extract_images and hasattr(shape, "image"):
                        try:
                            images.append(
                                {
                                    "slide": slide_idx + 1,
                                    "name": getattr(
                                        shape.image, "filename", f"image_{len(images)}"
                                    ),
                                    "size": (
                                        (shape.width, shape.height)
                                        if hasattr(shape, "width")
                                        else None
                                    ),
                                }
                            )
                        except Exception as e:
                            logging.warning(
                                f"Failed to extract image from slide {slide_idx + 1}: {str(e)}"
                            )

                except Exception as e:
                    logging.warning(
                        f"Failed to process shape on slide {slide_idx + 1}: {str(e)}"
                    )
                    continue

            # Extract hyperlinks if enabled
            if config.extract_links:
                try:
                    for shape in slide.shapes:
                        if hasattr(shape, "click_action") and shape.click_action:
                            if (
                                hasattr(shape.click_action, "hyperlink")
                                and shape.click_action.hyperlink
                            ):
                                hyperlink = shape.click_action.hyperlink
                                if hasattr(hyperlink, "address") and hyperlink.address:
                                    links.append(
                                        {
                                            "slide": slide_idx + 1,
                                            "text": getattr(shape, "text", ""),
                                            "url": hyperlink.address,
                                        }
                                    )
                except Exception as e:
                    logging.warning(
                        f"Failed to extract links from slide {slide_idx + 1}: {str(e)}"
                    )

            # Add slide content
            if slide_content:
                slide_text = f"=== Slide {slide_idx + 1} ===\n" + "\n".join(
                    slide_content
                )
                content_parts.append(slide_text)

        # Extract speaker notes
        for slide_idx, slide in enumerate(prs.slides):
            try:
                if hasattr(slide, "notes_slide") and slide.notes_slide:
                    notes_text = ""
                    for shape in slide.notes_slide.shapes:
                        if hasattr(shape, "text"):
                            notes_text += shape.text

                    if notes_text.strip():
                        notes_content = f"Speaker Notes (Slide {slide_idx + 1}): {notes_text.strip()}"
                        content_blocks.append(
                            ContentBlock(
                                content=notes_content,
                                block_type="notes",
                                page_number=slide_idx + 1,
                            )
                        )
                        content_parts.append(notes_content)
            except Exception as e:
                logging.warning(
                    f"Failed to extract notes from slide {slide_idx + 1}: {str(e)}"
                )

        # Combine content
        document.content = "\n\n".join(content_parts)
        document.content_blocks = content_blocks
        document.tables = tables
        document.images = images
        document.links = links

        # Update metadata
        document.metadata.word_count = len(document.content.split())
        document.metadata.character_count = len(document.content)

        return document
