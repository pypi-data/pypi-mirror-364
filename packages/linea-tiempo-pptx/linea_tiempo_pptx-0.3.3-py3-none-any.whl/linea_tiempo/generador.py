import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_LINE_DASH_STYLE
import calendar
import os
from collections import defaultdict

def generar_presentacion(excel_path, output_path="linea_tiempo_output.pptx", 
                         ordenamiento=None, graficos_por_diapositiva=2):
    
    # Validar número de gráficos por diapositiva
    if graficos_por_diapositiva not in [1, 2]:
        print("¡Advertencia! Solo se permiten 1 o 2 gráficos por diapositiva. Se usará 2 por defecto.")
        graficos_por_diapositiva = 2

    # Verificar si el archivo existe
    if not os.path.exists(excel_path):
        raise FileNotFoundError(f"No se encontró el archivo Excel en la ruta: {excel_path}")

    # Cargar archivo Excel
    df = pd.read_excel(excel_path, engine="openpyxl")

    # Convertir columnas de fecha a tipo datetime
    df['Fecha'] = pd.to_datetime(df['Fecha'])
    
    # Mantener orden original si no se especifica ordenamiento
    if ordenamiento:
        # Verificar columnas existentes
        columnas_validas = [col for col in ordenamiento if col in df.columns]
        columnas_no_validas = set(ordenamiento) - set(columnas_validas)
        
        if columnas_no_validas:
            print(f"Advertencia: Columnas de ordenamiento no encontradas: {columnas_no_validas}")
        
        # Ordenar solo si hay columnas válidas
        if columnas_validas:
            df = df.sort_values(by=columnas_validas + ['Iniciativa','Fecha'])    
   
    iniciativas = df['Iniciativa'].unique()

    # Crear presentación
    prs = Presentation()
    prs.slide_width = Cm(33.887)
    prs.slide_height = Cm(19.05)
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Colores para los estados de los hitos
    color_estado = {
        'Cumplido': RGBColor(0, 176, 80),
        'En gestión': RGBColor(255, 192, 0),
        'Proyección': RGBColor(255, 0, 0)
    }

    # Función para generar una escala de tiempo suave
    def generar_escala_proporcion_suave(fechas):
        fechas = sorted(fechas)
        escala = []
        for i in range(len(fechas) - 1):
            inicio = fechas[i]
            fin = fechas[i + 1]
            escala.append(inicio)
            meses_diff = (fin.year - inicio.year) * 12 + (fin.month - inicio.month)
            if meses_diff > 6:
                pasos = min(3, meses_diff - 1)
                for j in range(1, pasos + 1):
                    nuevo_mes = (inicio.to_timestamp() + pd.DateOffset(months=j)).to_period('M')
                    escala.append(nuevo_mes)
        escala.append(fechas[-1])
        return sorted(set(escala), key=lambda x: x.start_time)

    # Procesar iniciativas según el número especificado por diapositiva
    for i in range(0, len(iniciativas), graficos_por_diapositiva):
        iniciativas_par = iniciativas[i:i+graficos_por_diapositiva]
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        df_par = df[df['Iniciativa'].isin(iniciativas_par)]

        fechas_hitos = df_par['Fecha'].dt.to_period('M').unique()
        meses_clave = {mes for mes in fechas_hitos if mes.month in [1, 6, 12]}
        escala_suave = generar_escala_proporcion_suave(fechas_hitos)
        meses_total = sorted(set(escala_suave).union(meses_clave), key=lambda x: x.start_time)

        left_margin = Inches(2.5)
        top_margin = Inches(2.0)
        width_total = slide_width - left_margin - Inches(0.5)

        # Dibujar líneas verticales para cada mes
        for idx, mes in enumerate(meses_total):
            x = left_margin + idx * (width_total / len(meses_total))
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, top_margin, Pt(1), Inches(5))
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor(200, 200, 200)
            line.line.fill.background()

        # Crear diccionario de meses en español
        meses_es = {
            1: 'Ene', 2: 'Feb', 3: 'Mar', 4: 'Abr', 5: 'May', 6: 'Jun',
            7: 'Jul', 8: 'Ago', 9: 'Sep', 10: 'Oct', 11: 'Nov', 12: 'Dic'
        }

        # Dibujar etiquetas de meses
        for idx, mes in enumerate(meses_total):
            x = left_margin + idx * (width_total / len(meses_total))        
            mes_nombre = meses_es[mes.month]
            label = f"{mes_nombre} {mes.year}" if mes.month in [1, 6, 12] else f"{mes_nombre}"

            textbox = slide.shapes.add_textbox(x, top_margin - Inches(0.3), Inches(1), Inches(0.3))
            tf = textbox.text_frame
            p = tf.paragraphs[0]
            p.text = label
            p.font.size = Pt(10)
            p.alignment = PP_ALIGN.CENTER

        # Ajustar posición vertical según número de gráficos
        if graficos_por_diapositiva == 1:
            # Centrar verticalmente para un solo gráfico
            y_offsets = [top_margin + Inches(2.0)]
        else:
            # Posiciones para dos gráficos
            y_offsets = [top_margin + Inches(1.0), top_margin + Inches(3.2)]

        # Dibujar líneas horizontales para cada iniciativa
        for j, iniciativa in enumerate(iniciativas_par):
            y_offset = y_offsets[j]

            # Título de la iniciativa
            title_box = slide.shapes.add_textbox(Inches(0.3), y_offset - Inches(0.2), Inches(2), Inches(0.5))
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = iniciativa
            p.font.size = Pt(12)
            p.alignment = PP_ALIGN.JUSTIFY
            tf.word_wrap = True

            # Línea horizontal para la iniciativa
            linea = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                left_margin,
                y_offset + Inches(0.1),
                left_margin + width_total,
                y_offset + Inches(0.1)
            )
            linea.line.color.rgb = RGBColor(0, 0, 0)
            linea.line.dash_style = MSO_LINE_DASH_STYLE.DASH
            linea.line.width = Pt(2)

            df_iniciativa = df[df['Iniciativa'] == iniciativa].reset_index(drop=True)
            fechas_dict = defaultdict(list)
            for k, row in df_iniciativa.iterrows():
                fechas_dict[row['Fecha'].date()].append((k, row))

            global_label_counter = 0

            # Procesar cada mes y sus hitos
            for mes in df_iniciativa['Fecha'].dt.to_period('M').unique():
                df_mes = df_iniciativa[df_iniciativa['Fecha'].dt.to_period('M') == mes]
                fechas_unicas = df_mes['Fecha'].dt.date.unique()
                idx_mes = meses_total.index(mes)
                x_start = left_margin + idx_mes * (width_total / len(meses_total))
                x_end = left_margin + (idx_mes + 1) * (width_total / len(meses_total))

                # Procesar cada fecha única en el mes
                for fecha in sorted(fechas_unicas):
                    rows_fecha = df_mes[df_mes['Fecha'].dt.date == fecha]
                    same_date = len(rows_fecha) > 1
                    for pos, (k, row) in enumerate(rows_fecha.iterrows()):
                        # Calcular posición X del hito
                        if same_date:
                            day = row['Fecha'].day
                            days_in_month = row['Fecha'].days_in_month
                            x = x_start + (day / days_in_month) * (x_end - x_start)
                            y_variation = Inches(0.25 * pos)
                            label_mode = pos
                        elif len(fechas_unicas) == 1:
                            day = row['Fecha'].day
                            days_in_month = row['Fecha'].days_in_month
                            x = x_start + (day / days_in_month) * (x_end - x_start) - Inches(0.1)
                            y_variation = Inches(0)
                            label_mode = global_label_counter % 2
                            global_label_counter += 1
                        else:
                            spacing = (x_end - x_start) / (len(fechas_unicas) + 1)
                            x = x_start + spacing * (list(fechas_unicas).index(fecha) + 1)
                            y_variation = Inches(0)
                            label_mode = global_label_counter % 2
                            global_label_counter += 1

                        estado = row['Estado']
                        color = color_estado.get(estado, RGBColor(128, 128, 128))

                        # Dibujar marcador del hito
                        rh = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, x, y_offset + y_variation, Inches(0.2), Inches(0.2))
                        rh.fill.solid()
                        rh.fill.fore_color.rgb = color
                        rh.line.color.rgb = RGBColor(0, 0, 0)

                         # Posicionamiento de etiquetas
                        if same_date:
                            if label_mode == 0:
                                label_top = y_offset + y_variation - Inches(0.65)
                                label_left = x - Inches(0.6)
                                alignment = PP_ALIGN.CENTER
                            elif label_mode % 2 == 1:
                                label_top = y_offset + y_variation
                                label_left = x + Inches(0.3)
                                alignment = PP_ALIGN.LEFT
                            else:
                                label_top = y_offset + y_variation
                                label_left = x - Inches(1.6)
                                alignment = PP_ALIGN.RIGHT
                        else:
                            if label_mode == 0:
                                label_top = y_offset + y_variation - Inches(0.65)
                            else:
                                label_top = y_offset + y_variation + Inches(0.3)
                            label_left = x - Inches(0.6)
                            alignment = PP_ALIGN.CENTER

                        # Etiqueta del hito
                        label_box = slide.shapes.add_textbox(label_left, label_top, Inches(1.4), Inches(0.6))
                        tf = label_box.text_frame
                        p = tf.paragraphs[0]
                        p.text = f"{row['Hito']}\n{row['Fecha'].strftime('%d.%m.%Y')}"
                        p.font.size = Pt(8)
                        p.alignment = alignment
                        tf.word_wrap = True

    # Guardar presentación
    prs.save(output_path)